"""
create_pptx.py — Convert a markdown analysis file into a PowerPoint presentation.

Usage:
    python create_pptx.py <input.md> [output.pptx]

Requires: python-pptx (pip install python-pptx)
"""

import sys
import os
import re

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx is not installed. Run: pip install python-pptx")
    sys.exit(1)


def parse_markdown_sections(md_path: str) -> list[dict]:
    """Parse a markdown file into sections based on ## headings."""
    if not os.path.isfile(md_path):
        print(f"Error: File not found: {md_path}")
        sys.exit(1)

    with open(md_path, "r", encoding="utf-8") as f:
        content = f.read()

    sections: list[dict] = []
    current_section: dict | None = None

    for line in content.split("\n"):
        heading_match = re.match(r"^##\s+(.+)", line)
        if heading_match:
            if current_section:
                sections.append(current_section)
            current_section = {
                "title": heading_match.group(1).strip(),
                "bullets": [],
            }
        elif current_section:
            stripped = line.strip()
            if stripped.startswith(("- ", "* ", "• ")):
                bullet_text = re.sub(r"^[-*•]\s+", "", stripped)
                current_section["bullets"].append(bullet_text)
            elif stripped.startswith(("_", "*")) and stripped.endswith(("_", "*")):
                # Italic "So what?" line
                clean = stripped.strip("_* ")
                current_section["bullets"].append(clean)

    if current_section:
        sections.append(current_section)

    return sections


def create_presentation(sections: list[dict], output_path: str) -> None:
    """Create a PowerPoint presentation from parsed sections."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title Slide ---
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Volvo Group"
    subtitle.text = "Q4 & Full Year 2025 — Executive Brief"

    # Style title
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(0x4A, 0x4A, 0x4A)

    # --- Content Slides ---
    for section in sections:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = section["title"]
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        bullets = section["bullets"][:6]  # Max 6 bullets per slide
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Check for warning flag
            has_warning = "⚠️" in bullet or "⚠" in bullet
            clean_text = bullet.replace("⚠️", "").replace("⚠", "").strip()

            p.text = clean_text
            p.font.size = Pt(18)
            p.space_after = Pt(8)

            if has_warning:
                p.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)  # Red for declines
            else:
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Closing Slide ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Key Takeaway"
    for paragraph in slide.shapes.title.text_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = (
        "Review the financial-analysis.md for the full narrative. "
        "This deck was auto-generated from that analysis."
    )
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0x4A, 0x4A, 0x4A)
    p.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    print(f"Done. Presentation saved to {output_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python create_pptx.py <input.md> [output.pptx]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else "executive-brief.pptx"

    sections = parse_markdown_sections(input_file)
    if not sections:
        print("Warning: No sections found in the input file.")
        print("Expected markdown with ## headings and bullet points.")
        sys.exit(1)

    create_presentation(sections, output_file)
